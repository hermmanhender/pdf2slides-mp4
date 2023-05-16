import os
from pptx import Presentation
import PyPDF2

def extract_text_from_pdf(file_path):
    """Extrae el texto del archivo PDF y lo devuelve como una cadena."""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfFileReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        return text

def create_presentation_from_text(text):
    """Crea una presentación en PowerPoint a partir del texto proporcionado."""
    presentation = Presentation()
    slides = text.split('\n\n')  # Divide el texto en diapositivas basadas en párrafos separados
    for slide_content in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = slide_content.split('\n')[0]  # Toma la primera línea del párrafo como título de la diapositiva
        content.text = slide_content  # El párrafo completo se agrega como contenido de la diapositiva
    return presentation

def save_presentation(presentation, output_path):
    """Guarda la presentación en el archivo especificado."""
    presentation.save(output_path)

def convert_pdf_to_presentation(pdf_file, output_file):
    """Convierte un archivo PDF en una presentación en PowerPoint."""
    text = extract_text_from_pdf(pdf_file)
    presentation = create_presentation_from_text(text)
    save_presentation(presentation, output_file)

# Ejemplo de uso:
pdf_file_path = 'ruta/al/archivo.pdf'
output_file_path = 'ruta/de/salida.pptx'
convert_pdf_to_presentation(pdf_file_path, output_file_path)
print("Presentación generada con éxito en", output_file_path)

import os
from pptx import Presentation
import PyPDF2
import openai

openai.api_key = 'TU_CLAVE_DE_API_DE_OPENAI'  # Reemplaza con tu clave de API de OpenAI

def extract_notes_and_structure_from_pdf(file_path):
    """Extrae las notas del PDF y devuelve una lista de títulos, subtítulos y notas."""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfFileReader(file)
        structure = []
        current_title = ''
        current_subtitle = ''
        current_notes = ''
        for page in reader.pages:
            text = page.extract_text()
            lines = text.split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    if line.isupper():
                        if current_title:
                            structure.append((current_title, current_subtitle, current_notes))
                            current_notes = ''
                        current_title = line
                    else:
                        if current_subtitle:
                            current_notes += line + '\n'
                        else:
                            current_subtitle = line
        if current_title:
            structure.append((current_title, current_subtitle, current_notes))
        return structure

def generate_slide_content(title, subtitle, notes):
    """Genera el contenido de una diapositiva en base al título, subtítulo y notas."""
    prompt = f"# {title}\n## {subtitle}\n{notes}"
    response = openai.Completion.create(
        engine='text-davinci-003',
        prompt=prompt,
        max_tokens=200,
        n=1,
        stop=None,
        temperature=0.7
    )
    slide_content = response.choices[0].text.strip()
    return slide_content

def create_presentation_from_structure(structure):
    """Crea una presentación en PowerPoint en base a la estructura proporcionada."""
    presentation = Presentation()
    for title, subtitle, notes in structure:
        slide_content = generate_slide_content(title, subtitle, notes)
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = slide_content
    return presentation

def save_presentation(presentation, output_path):
    """Guarda la presentación en el archivo especificado."""
    presentation.save(output_path)

def convert_pdf_to_presentation(pdf_file, output_file):
    """Convierte un archivo PDF en una presentación en PowerPoint."""
    structure = extract_notes_and_structure_from_pdf(pdf_file)
    presentation = create_presentation_from_structure(structure)
    save_presentation(presentation, output_file)

# Ejemplo de uso:
pdf_file_path = 'ruta/al/archivo.pdf'
output_file_path = 'ruta/de/salida.pptx'
convert_pdf_to_presentation(pdf_file_path, output_file_path)
print("Presentación generada con éxito en", output_file_path)


from pptx2pdf import convert
from moviepy.editor import *
from gtts import gTTS
from pydub import AudioSegment

def convert_text_to_audio(text, output_file):
    """Convierte el texto en audio y lo guarda en el archivo especificado."""
    tts = gTTS(text=text, lang='es')
    tts.save(output_file)

def convert_presentation_to_video(presentation_file, audio_files, output_file):
    """Convierte una presentación en PowerPoint en un video con audio y subtítulos."""
    # Convierte la presentación a PDF
    pdf_file = 'temp.pdf'
    convert(presentation_file, pdf_file)

    # Carga las diapositivas en PDF como clips de imágenes
    clips = [ImageClip(pdf_file, duration=5) for _ in range(len(audio_files))]

    # Combina cada diapositiva con su respectivo audio
    for i, clip in enumerate(clips):
        audio = AudioSegment.from_file(audio_files[i])
        audio = audio.set_frame_rate(44100)  # Ajusta la frecuencia de muestreo si es necesario
        clip = clip.set_audio(audio)

        # Agrega subtítulos a la diapositiva
        txt_clip = TextClip(audio_files[i], fontsize=24, color='white', bg_color='black', method='caption').set_duration(5)
        txt_clip = txt_clip.set_position(('center', 'bottom')).set_audio(audio)
        clip = CompositeVideoClip([clip, txt_clip])

        clips[i] = clip

    # Crea el video final a partir de los clips de imágenes combinados
    video_clip = concatenate_videoclips(clips)

    # Ajusta el tamaño del video y otros parámetros
    video_clip = video_clip.resize(height=720)  # Cambia la altura del video a 720p (puedes ajustarlo según tus necesidades)

    # Guarda el video en el formato deseado
    video_clip.write_videofile(output_file, codec='libx264', audio_codec='aac')

# Ejemplo de uso:
presentation_file = 'ruta/de/la/presentacion.pptx'
notes = ['Texto de las notas 1', 'Texto de las notas 2', 'Texto de las notas 3']
audio_files = ['audio1.mp3', 'audio2.mp3', 'audio3.mp3']
output_video_file = 'ruta/del/video.mp4'

for i, note in enumerate(notes):
    audio_file = f'audio{i + 1}.mp3'
    convert_text_to_audio(note, audio_file)

convert_presentation_to_video(presentation_file, audio_files, output_video_file)
print("Video generado con éxito en", output_video_file)
